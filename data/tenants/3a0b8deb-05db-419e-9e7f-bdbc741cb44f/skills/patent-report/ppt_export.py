"""Export sections to PowerPoint."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from io_utils import file_artifact, sanitize_filename
from pptx import Presentation
from pptx.util import Inches, Pt


def export_ppt(
    sections: list[dict[str, Any]],
    *,
    out_dir: Path,
    key_prefix: str,
    title: str,
    filename: str,
    charts: list[dict[str, Any]] | None = None,
    summary: str = "",
    interpretation: str = "",
) -> dict[str, Any]:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if title_slide.placeholders[1].text_frame:
        subtitle = summary[:200] if summary else f"{len(sections)} 个分析维度"
        title_slide.placeholders[1].text = subtitle

    if summary or interpretation:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "总结与解读"
        body = slide.placeholders[1].text_frame
        body.clear()
        if summary:
            p = body.paragraphs[0]
            p.text = f"总结：{summary}"
            p.font.size = Pt(14)
        if interpretation:
            p2 = body.add_paragraph()
            p2.text = f"解读：{interpretation}"
            p2.font.size = Pt(12)

    chart_by_name = {str(c.get("name") or ""): c for c in charts or []}

    for section in sections:
        section_title = str(section.get("title") or "Section")
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section_title
        body = slide.placeholders[1].text_frame
        body.clear()

        rows = section.get("rows") or []
        if rows:
            for item in rows[:12]:
                p = body.add_paragraph()
                p.text = f"{item.get('label')}: {item.get('value')}"
                p.level = 0
                p.font.size = Pt(16)
        elif section.get("text"):
            p = body.paragraphs[0]
            p.text = str(section.get("text"))[:1200]

        chart = chart_by_name.get(section_title)
        if chart and chart.get("path"):
            slide.shapes.add_picture(
                str(chart["path"]),
                Inches(6.2),
                Inches(1.2),
                width=Inches(3.2),
            )

    fname = sanitize_filename(filename, default="patent-report") + ".pptx"
    path = out_dir / fname
    prs.save(path)
    key = f"{key_prefix}/{fname}"
    return file_artifact(path, key, fmt="pptx")
